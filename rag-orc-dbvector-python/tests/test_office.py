"""Office document loading: PowerPoint (slide=page) and Excel (sheet=page),
plus the aggregate supported-file check. No DB / model / network needed.
"""

from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

from src.ingest.document_loader import (
    csv_to_pages,
    docx_to_pages,
    is_supported_file,
    pptx_to_pages,
    xlsx_to_pages,
)


def make_pptx(path: Path, slides: list[str]) -> Path:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text in slides:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        box.text_frame.text = text
    prs.save(str(path))
    return path


def make_xlsx(path: Path, sheets: dict[str, list[list]]) -> Path:
    wb = Workbook()
    wb.remove(wb.active)  # drop the default sheet
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    wb.save(str(path))
    return path


def test_pptx_one_slide_per_page(tmp_path):
    p = make_pptx(tmp_path / "deck.pptx", ["Slide one about PPE", "Slide two about grounding"])
    pages = pptx_to_pages(p)
    assert len(pages) == 2  # 1 slide = 1 page
    assert "PPE" in pages[0]
    assert "grounding" in pages[1]


def test_xlsx_one_sheet_per_page(tmp_path):
    p = make_xlsx(
        tmp_path / "book.xlsx",
        {
            "Params": [["name", "limit"], ["voltage", 110]],
            "Notes": [["check", "before energizing"]],
        },
    )
    pages = xlsx_to_pages(p)
    assert len(pages) == 2  # 1 sheet = 1 page
    assert "[Sheet: Params]" in pages[0]
    assert "voltage" in pages[0] and "110" in pages[0]  # numbers stringified
    assert "before energizing" in pages[1]


def make_docx(path: Path, paragraphs: list[str], table: list[list] | None = None) -> Path:
    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    if table:
        t = doc.add_table(rows=len(table), cols=len(table[0]))
        for i, row in enumerate(table):
            for j, val in enumerate(row):
                t.cell(i, j).text = str(val)
    doc.save(str(path))
    return path


def test_docx_paragraphs_and_tables(tmp_path):
    p = make_docx(
        tmp_path / "doc.docx",
        ["Manual de exploatare dulap LCC.", "Verificati protectiile inainte de punere sub tensiune."],
        table=[["parametru", "valoare"], ["tensiune", "110 kV"]],
    )
    text = "\n".join(docx_to_pages(p, block_size=1500))
    assert "Manual de exploatare" in text
    assert "protectiile" in text
    assert "110 kV" in text  # table cell captured


def test_docx_paginates_long_document(tmp_path):
    p = make_docx(
        tmp_path / "long.docx",
        [f"Paragraful {i} cu ceva continut tehnic relevant despre statii." for i in range(200)],
    )
    pages = docx_to_pages(p, block_size=500)
    assert len(pages) > 1  # a long doc is split into multiple retrievable blocks


def test_csv_to_pages_sniffs_delimiter(tmp_path):
    p = tmp_path / "data.csv"
    p.write_text("param;valoare;unitate\ncurent;1250;A\ntensiune;110;kV\n", encoding="utf-8")
    text = "\n".join(csv_to_pages(p, block_size=1500))
    assert "curent" in text and "1250" in text
    assert "\t" in text  # ';' delimiter normalized to tab-separated cells


def test_doc_legacy_word(tmp_path):
    if not shutil.which("textutil"):
        pytest.skip("no .doc converter (textutil) available")
    from src.ingest.document_loader import doc_to_pages

    txt = tmp_path / "src.txt"
    txt.write_text("Regulament mentenanta preventiva.\nVerificati protectiile inainte de lucru.\n")
    doc = tmp_path / "out.doc"
    subprocess.run(["textutil", "-convert", "doc", "-output", str(doc), str(txt)], check=True)
    text = "\n".join(doc_to_pages(doc, block_size=1500))
    assert "Regulament" in text or "mentenanta" in text


def test_is_supported_file_aggregate():
    for ok in ["a.pdf", "A.PDF", "x.png", "y.JPG", "p.pptx", "e.xlsx", "f.XLSM",
               "g.xls", "h.docx", "j.doc", "i.csv"]:
        assert is_supported_file(ok) is True, ok
    # legacy binary PowerPoint, videos, archives, and other types are not supported
    for no in ["a.txt", "c.ppt", "d.mov", "e.rar", "g"]:
        assert is_supported_file(no) is False, no
