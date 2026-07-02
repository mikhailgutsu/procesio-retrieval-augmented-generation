"""Unified document loading: turn any supported input into per-page text.

Dispatch by file type:
  * PDF / image        → text-layer PDF (OCR when needed) via `pdf_loader`;
                         `text_pdf_path` is set so spans can be highlighted on it.
  * PowerPoint (.pptx) → 1 slide = 1 page.
  * Excel (.xlsx/.xlsm)→ 1 sheet = 1 page.

Office formats have no source PDF, so `text_pdf_path` is None and highlighting
falls back to the UI payload (character offsets) — no PDF annotation.
"""

from __future__ import annotations

import csv
import io
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path

from ..config import Settings, get_settings
from ..errors import RagError
from ..logging_config import get_logger
from .pdf_loader import (
    IMAGE_SUFFIXES,
    ensure_text_pdf,
    extract_pages_text,
    render_page_png,
    sanitize_text,
    vision_transcribe_page,
)

log = get_logger(__name__)

PDF_SUFFIXES = {".pdf"}
PPTX_SUFFIXES = {".pptx"}           # python-pptx reads .pptx only (legacy .ppt unsupported)
XLSX_SUFFIXES = {".xlsx", ".xlsm"}  # openpyxl (modern XML spreadsheets)
XLS_SUFFIXES = {".xls"}             # xlrd (legacy binary spreadsheets)
DOCX_SUFFIXES = {".docx"}           # python-docx (modern XML Word)
DOC_SUFFIXES = {".doc"}             # legacy binary Word via textutil/antiword/soffice
CSV_SUFFIXES = {".csv"}
OFFICE_SUFFIXES = (
    PPTX_SUFFIXES | XLSX_SUFFIXES | XLS_SUFFIXES | DOCX_SUFFIXES | DOC_SUFFIXES | CSV_SUFFIXES
)
SUPPORTED_SUFFIXES = PDF_SUFFIXES | IMAGE_SUFFIXES | OFFICE_SUFFIXES


def is_supported_file(path: str | Path) -> bool:
    """True for any ingestable type: PDF, image, PowerPoint (.pptx),
    Excel (.xlsx/.xlsm/.xls), Word (.docx/.doc), CSV."""
    return Path(path).suffix.lower() in SUPPORTED_SUFFIXES


def _paginate_lines(lines: list[str], block_size: int) -> list[str]:
    """Group text lines into ~block_size-char pages, breaking at line boundaries."""
    pages: list[str] = []
    cur: list[str] = []
    n = 0
    for ln in lines:
        if n and n + len(ln) > block_size:
            pages.append("\n".join(cur))
            cur, n = [], 0
        cur.append(ln)
        n += len(ln) + 1
    if cur:
        pages.append("\n".join(cur))
    return pages


@dataclass
class LoadedDocument:
    pages: list[str]           # per-page text (index 0 == page/slide/sheet 1)
    text_pdf_path: str | None  # source PDF for highlighting; None for office formats
    was_ocred: bool
    kind: str                  # "pdf" | "image" | "pptx" | "xlsx"


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint — 1 slide = 1 page
# ─────────────────────────────────────────────────────────────────────────────
def _iter_shapes(shapes):
    """Yield leaf shapes, recursing into grouped shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def pptx_to_pages(path: str | Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    if any(cell.strip() for cell in cells):
                        parts.append("\t".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append("Note: " + notes)
        pages.append(sanitize_text("\n".join(parts)))
    return pages


# ─────────────────────────────────────────────────────────────────────────────
# Excel — 1 sheet = 1 page
# ─────────────────────────────────────────────────────────────────────────────
def xlsx_to_pages(path: str | Path) -> list[str]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        pages: list[str] = []
        for ws in wb.worksheets:
            lines = [f"[Sheet: {ws.title}]"]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
            pages.append(sanitize_text("\n".join(lines)))
        return pages
    finally:
        wb.close()


def xls_to_pages(path: str | Path) -> list[str]:
    """Legacy binary Excel (.xls) — 1 sheet = 1 page (via xlrd)."""
    import xlrd

    wb = xlrd.open_workbook(str(path))
    pages: list[str] = []
    for sh in wb.sheets():
        lines = [f"[Sheet: {sh.name}]"]
        for r in range(sh.nrows):
            cells = [str(sh.cell_value(r, c)) for c in range(sh.ncols) if str(sh.cell_value(r, c)).strip()]
            if cells:
                lines.append("\t".join(cells))
        pages.append(sanitize_text("\n".join(lines)))
    return pages or [""]


# ─────────────────────────────────────────────────────────────────────────────
# Word (.docx) — paragraphs + tables, paginated into ~block_size-char blocks
# (a .docx has no fixed pages; each block becomes one retrievable "page").
# ─────────────────────────────────────────────────────────────────────────────
def docx_to_pages(path: str | Path, block_size: int = 1500) -> list[str]:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(str(path))
    lines: list[str] = []
    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            text = Paragraph(child, document).text.strip()
            if text:
                lines.append(text)
        elif child.tag == qn("w:tbl"):
            for row in Table(child, document).rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    lines.append("\t".join(cells))
    pages = _paginate_lines(lines, block_size)
    return [sanitize_text(p) for p in pages] if pages else [""]


# ─────────────────────────────────────────────────────────────────────────────
# Legacy Word (.doc) — extract plain text with an external converter, then paginate.
# python-docx cannot read .doc; try macOS `textutil`, then `antiword`, then LibreOffice.
# ─────────────────────────────────────────────────────────────────────────────
def _doc_extract_text(path: Path) -> str:
    src = str(path)
    if shutil.which("textutil"):  # macOS built-in
        out = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", src], capture_output=True, timeout=180
        )
        if out.returncode == 0 and out.stdout.strip():
            return out.stdout.decode("utf-8", errors="replace")
    if shutil.which("antiword"):
        out = subprocess.run(["antiword", src], capture_output=True, timeout=180)
        if out.returncode == 0 and out.stdout.strip():
            return out.stdout.decode("utf-8", errors="replace")
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        with tempfile.TemporaryDirectory() as td:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", td, src],
                capture_output=True,
                timeout=240,
            )
            txt = Path(td) / (path.stem + ".txt")
            if txt.exists():
                return txt.read_text(encoding="utf-8", errors="replace")
    raise RagError(
        f"Cannot read legacy .doc {path.name}: no converter found. Install one of "
        "`antiword` or LibreOffice (macOS has `textutil` built-in), or convert to .docx."
    )


def doc_to_pages(path: str | Path, block_size: int = 1500) -> list[str]:
    text = _doc_extract_text(Path(path))
    lines = [ln for ln in text.splitlines() if ln.strip()]
    pages = _paginate_lines(lines, block_size)
    return [sanitize_text(p) for p in pages] if pages else [""]


# ─────────────────────────────────────────────────────────────────────────────
# CSV — delimiter sniffed, paginated into ~block_size-char blocks
# ─────────────────────────────────────────────────────────────────────────────
def csv_to_pages(path: str | Path, block_size: int = 1500) -> list[str]:
    raw = Path(path).read_bytes()
    text = None
    for enc in ("utf-8-sig", "utf-8", "cp1250", "latin-1"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        text = raw.decode("latin-1", errors="replace")
    try:
        dialect = csv.Sniffer().sniff(text[:4096], delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel
    reader = csv.reader(io.StringIO(text), dialect)
    lines = ["\t".join(str(c) for c in row) for row in reader if any(str(c).strip() for c in row)]
    pages = _paginate_lines(lines, block_size)
    return [sanitize_text(p) for p in pages] if pages else [""]


# ─────────────────────────────────────────────────────────────────────────────
# Optional vision fallback for PDF/image pages still empty after OCR
# ─────────────────────────────────────────────────────────────────────────────
def _maybe_vision_fill(
    text_pdf: Path,
    pages: list[str],
    was_ocred: bool,
    settings: Settings,
    *,
    force_all: bool = False,
) -> list[str]:
    """Transcribe pages with the vision model when OCR text is missing or unreliable.

    Normally only pages whose OCR text is below ``scanned_char_threshold`` are sent
    to the vision model. For standalone images (``force_all=True``) OCR is unreliable
    even when it returns *some* (garbled) text, so every page is transcribed and the
    richer of {vision, OCR} is kept.
    """
    if not settings.ocr_vision_fallback:
        return pages
    if not (was_ocred or force_all):
        return pages
    filled = list(pages)
    for i, text in enumerate(pages):
        needs_vision = force_all or len(text.strip()) < settings.scanned_char_threshold
        if not needs_vision:
            continue
        log.info("Vision transcription on page %d of %s", i + 1, text_pdf.name)
        png = render_page_png(text_pdf, i)
        transcribed = vision_transcribe_page(png, settings)
        if transcribed and len(transcribed.strip()) >= len(text.strip()):
            filled[i] = transcribed
        elif transcribed and text.strip():
            # Keep both when OCR had unique content the model may have missed.
            filled[i] = f"{transcribed}\n\n{text}"
        elif transcribed:
            filled[i] = transcribed
    return filled


# ─────────────────────────────────────────────────────────────────────────────
# Dispatch
# ─────────────────────────────────────────────────────────────────────────────
def load_document(path: str | Path, settings: Settings | None = None) -> LoadedDocument:
    """Load any supported document into per-page text (+ optional source PDF)."""
    settings = settings or get_settings()
    path = Path(path)
    suffix = path.suffix.lower()

    if suffix in PPTX_SUFFIXES:
        pages = pptx_to_pages(path)
        log.info("Loaded PowerPoint %s — %d slide(s).", path.name, len(pages))
        return LoadedDocument(pages=pages, text_pdf_path=None, was_ocred=False, kind="pptx")

    if suffix in XLSX_SUFFIXES:
        pages = xlsx_to_pages(path)
        log.info("Loaded Excel %s — %d sheet(s).", path.name, len(pages))
        return LoadedDocument(pages=pages, text_pdf_path=None, was_ocred=False, kind="xlsx")

    if suffix in XLS_SUFFIXES:
        pages = xls_to_pages(path)
        log.info("Loaded Excel (legacy) %s — %d sheet(s).", path.name, len(pages))
        return LoadedDocument(pages=pages, text_pdf_path=None, was_ocred=False, kind="xls")

    if suffix in DOCX_SUFFIXES:
        pages = docx_to_pages(path, settings.chunk_window_size)
        log.info("Loaded Word %s — %d block(s).", path.name, len(pages))
        return LoadedDocument(pages=pages, text_pdf_path=None, was_ocred=False, kind="docx")

    if suffix in DOC_SUFFIXES:
        pages = doc_to_pages(path, settings.chunk_window_size)
        log.info("Loaded Word (legacy .doc) %s — %d block(s).", path.name, len(pages))
        return LoadedDocument(pages=pages, text_pdf_path=None, was_ocred=False, kind="doc")

    if suffix in CSV_SUFFIXES:
        pages = csv_to_pages(path, settings.chunk_window_size)
        log.info("Loaded CSV %s — %d block(s).", path.name, len(pages))
        return LoadedDocument(pages=pages, text_pdf_path=None, was_ocred=False, kind="csv")

    # PDF or image: route through the text-layer-PDF path (OCR when needed).
    text_pdf, was_ocred = ensure_text_pdf(path, settings)
    pages = extract_pages_text(text_pdf)
    is_image = suffix in IMAGE_SUFFIXES
    # For standalone images OCR is unreliable (garbled brand names, diagrams); always
    # let the vision model read them when the fallback is enabled.
    pages = _maybe_vision_fill(text_pdf, pages, was_ocred, settings, force_all=is_image)
    kind = "image" if is_image else "pdf"
    return LoadedDocument(
        pages=pages, text_pdf_path=str(text_pdf), was_ocred=was_ocred, kind=kind
    )
