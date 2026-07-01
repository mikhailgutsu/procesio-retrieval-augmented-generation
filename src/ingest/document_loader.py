"""Unified document loading: turn any supported input into per-page text.

Dispatch by file type:
  * PDF / image        → text-layer PDF (OCR when needed) via `pdf_loader`;
                         `text_pdf_path` is set so spans can be highlighted on it.
  * PowerPoint (.pptx) → 1 slide = 1 page.
  * Excel (.xlsx/.xlsm)→ 1 sheet = 1 page.

Office formats have no source PDF, so `text_pdf_path` is None and highlighting
falls back to the UI payload (character offsets) — no PDF annotation.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from ..config import Settings, get_settings
from ..logging_config import get_logger
from .pdf_loader import (
    IMAGE_SUFFIXES,
    ensure_text_pdf,
    extract_pages_text,
    render_page_png,
    sanitize_text,
    vision_transcribe_page,
)

log = get_logger(__name__)

PDF_SUFFIXES = {".pdf"}
PPTX_SUFFIXES = {".pptx"}          # python-pptx reads .pptx only (legacy .ppt unsupported)
XLSX_SUFFIXES = {".xlsx", ".xlsm"}  # openpyxl reads .xlsx/.xlsm only (legacy .xls unsupported)
OFFICE_SUFFIXES = PPTX_SUFFIXES | XLSX_SUFFIXES
SUPPORTED_SUFFIXES = PDF_SUFFIXES | IMAGE_SUFFIXES | OFFICE_SUFFIXES


def is_supported_file(path: str | Path) -> bool:
    """True for any ingestable type: PDF, image, PowerPoint (.pptx), Excel (.xlsx/.xlsm)."""
    return Path(path).suffix.lower() in SUPPORTED_SUFFIXES


@dataclass
class LoadedDocument:
    pages: list[str]           # per-page text (index 0 == page/slide/sheet 1)
    text_pdf_path: str | None  # source PDF for highlighting; None for office formats
    was_ocred: bool
    kind: str                  # "pdf" | "image" | "pptx" | "xlsx"


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint — 1 slide = 1 page
# ─────────────────────────────────────────────────────────────────────────────
def _iter_shapes(shapes):
    """Yield leaf shapes, recursing into grouped shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def pptx_to_pages(path: str | Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    if any(cell.strip() for cell in cells):
                        parts.append("\t".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append("Note: " + notes)
        pages.append(sanitize_text("\n".join(parts)))
    return pages


# ─────────────────────────────────────────────────────────────────────────────
# Excel — 1 sheet = 1 page
# ─────────────────────────────────────────────────────────────────────────────
def xlsx_to_pages(path: str | Path) -> list[str]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        pages: list[str] = []
        for ws in wb.worksheets:
            lines = [f"[Sheet: {ws.title}]"]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
            pages.append(sanitize_text("\n".join(lines)))
        return pages
    finally:
        wb.close()


# ─────────────────────────────────────────────────────────────────────────────
# Optional vision fallback for PDF/image pages still empty after OCR
# ─────────────────────────────────────────────────────────────────────────────
def _maybe_vision_fill(
    text_pdf: Path, pages: list[str], was_ocred: bool, settings: Settings
) -> list[str]:
    if not (was_ocred and settings.ocr_vision_fallback):
        return pages
    filled = list(pages)
    for i, text in enumerate(pages):
        if len(text.strip()) < settings.scanned_char_threshold:
            log.info("Vision fallback on page %d of %s", i + 1, text_pdf.name)
            png = render_page_png(text_pdf, i)
            transcribed = vision_transcribe_page(png, settings)
            if transcribed:
                filled[i] = transcribed
    return filled


# ─────────────────────────────────────────────────────────────────────────────
# Dispatch
# ─────────────────────────────────────────────────────────────────────────────
def load_document(path: str | Path, settings: Settings | None = None) -> LoadedDocument:
    """Load any supported document into per-page text (+ optional source PDF)."""
    settings = settings or get_settings()
    path = Path(path)
    suffix = path.suffix.lower()

    if suffix in PPTX_SUFFIXES:
        pages = pptx_to_pages(path)
        log.info("Loaded PowerPoint %s — %d slide(s).", path.name, len(pages))
        return LoadedDocument(pages=pages, text_pdf_path=None, was_ocred=False, kind="pptx")

    if suffix in XLSX_SUFFIXES:
        pages = xlsx_to_pages(path)
        log.info("Loaded Excel %s — %d sheet(s).", path.name, len(pages))
        return LoadedDocument(pages=pages, text_pdf_path=None, was_ocred=False, kind="xlsx")

    # PDF or image: route through the text-layer-PDF path (OCR when needed).
    text_pdf, was_ocred = ensure_text_pdf(path, settings)
    pages = extract_pages_text(text_pdf)
    pages = _maybe_vision_fill(text_pdf, pages, was_ocred, settings)
    kind = "image" if suffix in IMAGE_SUFFIXES else "pdf"
    return LoadedDocument(
        pages=pages, text_pdf_path=str(text_pdf), was_ocred=was_ocred, kind=kind
    )
