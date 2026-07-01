"""Office document loading: PowerPoint (slide=page) and Excel (sheet=page),
plus the aggregate supported-file check. No DB / model / network needed.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

from src.ingest.document_loader import is_supported_file, pptx_to_pages, xlsx_to_pages


def make_pptx(path: Path, slides: list[str]) -> Path:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text in slides:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        box.text_frame.text = text
    prs.save(str(path))
    return path


def make_xlsx(path: Path, sheets: dict[str, list[list]]) -> Path:
    wb = Workbook()
    wb.remove(wb.active)  # drop the default sheet
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    wb.save(str(path))
    return path


def test_pptx_one_slide_per_page(tmp_path):
    p = make_pptx(tmp_path / "deck.pptx", ["Slide one about PPE", "Slide two about grounding"])
    pages = pptx_to_pages(p)
    assert len(pages) == 2  # 1 slide = 1 page
    assert "PPE" in pages[0]
    assert "grounding" in pages[1]


def test_xlsx_one_sheet_per_page(tmp_path):
    p = make_xlsx(
        tmp_path / "book.xlsx",
        {
            "Params": [["name", "limit"], ["voltage", 110]],
            "Notes": [["check", "before energizing"]],
        },
    )
    pages = xlsx_to_pages(p)
    assert len(pages) == 2  # 1 sheet = 1 page
    assert "[Sheet: Params]" in pages[0]
    assert "voltage" in pages[0] and "110" in pages[0]  # numbers stringified
    assert "before energizing" in pages[1]


def test_is_supported_file_aggregate():
    for ok in ["a.pdf", "A.PDF", "x.png", "y.JPG", "d.pptx", "e.xlsx", "f.XLSM"]:
        assert is_supported_file(ok) is True, ok
    # legacy binary Office formats and other types are not supported
    for no in ["a.txt", "b.docx", "c.ppt", "d.xls", "g", "h.csv"]:
        assert is_supported_file(no) is False, no
